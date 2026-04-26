from flask import Flask
from threading import Thread

app = Flask('')

@app.route('/')
def home():
    return "Bot tirik!"

def run():
    app.run(host='0.0.0.0', port=8080)

def keep_alive():
    t = Thread(target=run)
    t.start()
import asyncio
import json
import logging
import os
import re
import sqlite3
import tempfile
import uuid
from collections import OrderedDict
from datetime import datetime
from io import BytesIO
from pathlib import Path

import aiofiles
import yt_dlp
from aiogram import Bot, Dispatcher, executor, types
from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.dispatcher import FSMContext
from aiogram.dispatcher.filters.state import State, StatesGroup
from aiogram.types import (
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    KeyboardButton,
    ReplyKeyboardMarkup,
)
from openai import AsyncOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

BOT_TOKEN = os.getenv("BOT_TOKEN")
OPENAI_BASE_URL = os.getenv("AI_INTEGRATIONS_OPENAI_BASE_URL")
OPENAI_API_KEY = os.getenv("AI_INTEGRATIONS_OPENAI_API_KEY")

if not BOT_TOKEN:
    raise RuntimeError("BOT_TOKEN environment variable is not set.")

bot = Bot(token=BOT_TOKEN, parse_mode="HTML")
storage = MemoryStorage()
dp = Dispatcher(bot, storage=storage)

ai_client = (
    AsyncOpenAI(base_url=OPENAI_BASE_URL, api_key=OPENAI_API_KEY)
    if OPENAI_BASE_URL and OPENAI_API_KEY
    else None
)

DOWNLOAD_DIR = Path(tempfile.gettempdir()) / "bot_downloads"
DOWNLOAD_DIR.mkdir(exist_ok=True)

CLOCK_FRAMES = ["🕐", "🕑", "🕒", "🕓", "🕔", "🕕", "🕖", "🕗", "🕘", "🕙", "🕚", "🕛"]

# ---- Configuration ----
CARD_NUMBER = "5614 6835 1440 1090"
CARD_HOLDER = "Nurdiyor Fayzullayev"
FREE_LIMIT = 1  # 1 marta bepul (har bir talaba ishi turi uchun)

# Pricing (so'm). (pages, price)
SLIDES_PRICING = [(10, 2_000), (15, 5_000), (20, 10_000)]
ESSAY_PRICING = [(5, 2_000), (10, 5_000), (20, 10_000), (30, 15_000)]
COURSEWORK_PRICING = [(20, 10_000), (30, 15_000), (40, 20_000), (50, 25_000)]

REQUEST_LABELS = {
    "essay": ("📝 Referat", ESSAY_PRICING),
    "coursework": ("📚 Kurs ishi", COURSEWORK_PRICING),
    "slides": ("📊 Slayd", SLIDES_PRICING),
}

DB_PATH = Path("bot.db")

# Cache: short-key → social URL (for "music after video" button)
URL_CACHE: "OrderedDict[str, str]" = OrderedDict()
URL_CACHE_MAX = 1000


def cache_url(url: str) -> str:
    key = uuid.uuid4().hex[:10]
    URL_CACHE[key] = url
    while len(URL_CACHE) > URL_CACHE_MAX:
        URL_CACHE.popitem(last=False)
    return key


# ---------- DB ----------
def db_init():
    con = sqlite3.connect(DB_PATH)
    con.executescript(
        """
        CREATE TABLE IF NOT EXISTS users (
            user_id INTEGER PRIMARY KEY,
            username TEXT,
            full_name TEXT,
            essay_used INTEGER DEFAULT 0,
            coursework_used INTEGER DEFAULT 0,
            slides_used INTEGER DEFAULT 0,
            is_admin INTEGER DEFAULT 0,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS requests (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            request_type TEXT,
            topic TEXT,
            pages INTEGER,
            price INTEGER,
            status TEXT DEFAULT 'awaiting_payment',
            receipt_file_id TEXT,
            admin_msg_id INTEGER,
            created_at TEXT DEFAULT CURRENT_TIMESTAMP,
            reviewed_at TEXT
        );
        """
    )
    con.commit()
    con.close()


def db():
    con = sqlite3.connect(DB_PATH)
    con.row_factory = sqlite3.Row
    return con


def ensure_user(user: types.User) -> sqlite3.Row:
    con = db()
    cur = con.cursor()
    cur.execute("SELECT * FROM users WHERE user_id=?", (user.id,))
    row = cur.fetchone()
    if row is None:
        cur.execute("SELECT COUNT(*) FROM users WHERE is_admin=1")
        is_admin = 1 if cur.fetchone()[0] == 0 else 0
        cur.execute(
            "INSERT INTO users (user_id, username, full_name, is_admin) VALUES (?,?,?,?)",
            (user.id, user.username or "", user.full_name, is_admin),
        )
        con.commit()
        cur.execute("SELECT * FROM users WHERE user_id=?", (user.id,))
        row = cur.fetchone()
    else:
        cur.execute(
            "UPDATE users SET username=?, full_name=? WHERE user_id=?",
            (user.username or "", user.full_name, user.id),
        )
        con.commit()
    con.close()
    return row


def is_admin(user_id: int) -> bool:
    con = db()
    row = con.execute(
        "SELECT is_admin FROM users WHERE user_id=?", (user_id,)
    ).fetchone()
    con.close()
    return bool(row and row["is_admin"])


def get_admin_id() -> int | None:
    con = db()
    row = con.execute("SELECT user_id FROM users WHERE is_admin=1 LIMIT 1").fetchone()
    con.close()
    return row["user_id"] if row else None


def get_usage_count(user_id: int, field: str) -> int:
    con = db()
    row = con.execute(
        f"SELECT {field} FROM users WHERE user_id=?", (user_id,)
    ).fetchone()
    con.close()
    return row[field] if row else 0


def increment_usage(user_id: int, field: str):
    con = db()
    con.execute(f"UPDATE users SET {field}={field}+1 WHERE user_id=?", (user_id,))
    con.commit()
    con.close()


def has_free_quota(user_id: int, request_type: str) -> bool:
    if is_admin(user_id):
        return True
    field = f"{request_type}_used"
    return get_usage_count(user_id, field) < FREE_LIMIT


def create_request(user_id: int, rtype: str, topic: str, pages: int, price: int) -> int:
    con = db()
    cur = con.cursor()
    cur.execute(
        "INSERT INTO requests (user_id, request_type, topic, pages, price) VALUES (?,?,?,?,?)",
        (user_id, rtype, topic, pages, price),
    )
    rid = cur.lastrowid
    con.commit()
    con.close()
    return rid


def get_request(rid: int) -> sqlite3.Row | None:
    con = db()
    row = con.execute("SELECT * FROM requests WHERE id=?", (rid,)).fetchone()
    con.close()
    return row


def update_request(rid: int, **fields):
    if not fields:
        return
    sets = ", ".join(f"{k}=?" for k in fields)
    vals = list(fields.values()) + [rid]
    con = db()
    con.execute(f"UPDATE requests SET {sets} WHERE id=?", vals)
    con.commit()
    con.close()


# ---------- FSM ----------
class Form(StatesGroup):
    waiting_essay_topic = State()
    waiting_coursework_topic = State()
    waiting_slides_topic = State()
    waiting_receipt = State()  # data: request_id
    chatting_ai = State()


# ---------- Keyboards ----------
def main_menu(user_id: int | None = None) -> ReplyKeyboardMarkup:
    kb = ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row(KeyboardButton("📝 Referat"), KeyboardButton("📚 Kurs ishi"))
    kb.row(KeyboardButton("📊 Slayd (PPTX)"), KeyboardButton("🎵 Musiqa/Video"))
    kb.row(KeyboardButton("🤖 AI yordamchi"), KeyboardButton("👤 Profil"))
    kb.row(KeyboardButton("ℹ️ Yordam"))
    if user_id and is_admin(user_id):
        kb.row(KeyboardButton("🛠 Admin panel"))
    return kb


def cancel_kb() -> ReplyKeyboardMarkup:
    kb = ReplyKeyboardMarkup(resize_keyboard=True)
    kb.add(KeyboardButton("❌ Bekor qilish"))
    return kb


def pages_kb(rtype: str) -> InlineKeyboardMarkup:
    pricing = REQUEST_LABELS[rtype][1]
    kb = InlineKeyboardMarkup(row_width=2)
    for pages, price in pricing:
        kb.insert(
            InlineKeyboardButton(
                f"{pages} bet — {price:,} so'm",
                callback_data=f"sz:{rtype}:{pages}",
            )
        )
    return kb


def confirm_payment_kb(rid: int) -> InlineKeyboardMarkup:
    kb = InlineKeyboardMarkup()
    kb.add(
        InlineKeyboardButton(
            "✅ Men to'ladim — chek yuborish", callback_data=f"pay:{rid}"
        )
    )
    kb.add(InlineKeyboardButton("❌ Bekor qilish", callback_data=f"cancel:{rid}"))
    return kb


def admin_review_kb(rid: int) -> InlineKeyboardMarkup:
    kb = InlineKeyboardMarkup()
    kb.add(
        InlineKeyboardButton("✅ Qabul qilish", callback_data=f"adm:approve:{rid}"),
        InlineKeyboardButton("❌ Rad etish", callback_data=f"adm:reject:{rid}"),
    )
    return kb


def music_after_video_kb(key: str) -> InlineKeyboardMarkup:
    kb = InlineKeyboardMarkup()
    kb.add(
        InlineKeyboardButton("🎵 Musiqasini yuklab olish", callback_data=f"audio:{key}")
    )
    return kb


# ---------- Loading animation ----------
async def show_loading(chat_id: int, text: str = "Tayyorlanmoqda") -> types.Message:
    msg = await bot.send_message(chat_id, f"{CLOCK_FRAMES[0]} {text}...")

    async def animate():
        idx = 0
        try:
            while True:
                await asyncio.sleep(0.6)
                idx = (idx + 1) % len(CLOCK_FRAMES)
                try:
                    await bot.edit_message_text(
                        f"{CLOCK_FRAMES[idx]} {text}...",
                        chat_id=msg.chat.id,
                        message_id=msg.message_id,
                    )
                except Exception:
                    pass
        except asyncio.CancelledError:
            pass

    msg._anim_task = asyncio.create_task(animate())  # type: ignore[attr-defined]
    return msg


async def stop_loading(msg: types.Message):
    task = getattr(msg, "_anim_task", None)
    if task:
        task.cancel()
    try:
        await bot.delete_message(msg.chat.id, msg.message_id)
    except Exception:
        pass


# ---------- AI ----------
async def ai_text(
    prompt: str, system: str = "Siz foydali yordamchisiz. O'zbek tilida javob bering."
) -> str:
    if not ai_client:
        return "AI sozlanmagan."
    try:
        resp = await ai_client.chat.completions.create(
            model="gpt-5.4",
            max_completion_tokens=8192,
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": prompt},
            ],
        )
        return resp.choices[0].message.content or "Javob bo'sh."
    except Exception as e:
        logger.exception("AI error")
        return f"AI xatosi: {e}"


# ---------- Document builders ----------
def build_docx(title: str, body: str) -> BytesIO:
    doc = Document()
    doc.add_heading(title, 0)
    for para in body.split("\n"):
        para = para.strip()
        if not para:
            continue
        if para.startswith("# "):
            doc.add_heading(para[2:], level=1)
        elif para.startswith("## "):
            doc.add_heading(para[3:], level=2)
        elif para.startswith("### "):
            doc.add_heading(para[4:], level=3)
        else:
            doc.add_paragraph(para)
    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio


def build_pptx(title: str, slides_content: list[dict]) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    if len(s.placeholders) > 1:
        s.placeholders[1].text = "AI tomonidan yaratilgan taqdimot"
    bullet_layout = prs.slide_layouts[1]
    for sl in slides_content:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = sl.get("title", "")
        body = s.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(sl.get("bullets", [])):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.font.size = Pt(20)
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


async def generate_slides_content(topic: str, n: int) -> list[dict]:
    raw = await ai_text(
        f"'{topic}' mavzusida {n} ta slayd uchun JSON: "
        '[{"title":"...","bullets":["...","..."]}, ...]. '
        "Har slaydda 3-5 bullet. Faqat JSON qaytaring."
    )
    m = re.search(r"\[.*\]", raw, re.DOTALL)
    if m:
        try:
            data = json.loads(m.group(0))
            if isinstance(data, list) and data:
                return data
        except Exception:
            pass
    return [{"title": f"Slayd {i + 1}", "bullets": [topic]} for i in range(n)]


async def generate_essay(topic: str, pages: int) -> str:
    return await ai_text(
        f"'{topic}' mavzusida ~{pages} sahifalik (taxminan {pages * 350} so'z) to'liq referat yozing. "
        "Tarkibi: Kirish, Asosiy qism (kichik bo'limlarga ajrating, '## Bo'lim nomi' bilan), "
        "Xulosa, Foydalanilgan adabiyotlar (kamida 5 ta manba). "
        "Har bir bo'lim batafsil va akademik uslubda yozilsin.",
        system="Siz akademik yozuvchi, o'zbek tilida sifatli referatlar yozasiz.",
    )


async def generate_coursework(topic: str, pages: int) -> str:
    return await ai_text(
        f"'{topic}' mavzusida ~{pages} sahifalik (taxminan {pages * 350} so'z) to'liq kurs ishi yozing. "
        "Tarkibi (har bir bo'limni '## Sarlavha' bilan boshlang): "
        "Annotatsiya, Mundarija, Kirish (3+ paragraf), I bob - Nazariy qism (3+ kichik bo'lim), "
        "II bob - Tahliliy qism (3+ kichik bo'lim), III bob - Amaliy qism, "
        "Xulosa, Foydalanilgan adabiyotlar (kamida 10 ta manba). "
        "Ilmiy uslubda batafsil yozing.",
        system="Siz akademik tadqiqotchi, o'zbek tilida ilmiy kurs ishlari yozasiz.",
    )


# ---------- Media ----------
SOCIAL_HOSTS = (
    "youtube.com",
    "youtu.be",
    "tiktok.com",
    "instagram.com",
    "facebook.com",
    "fb.watch",
    "twitter.com",
    "x.com",
    "vk.com",
    "soundcloud.com",
    "vimeo.com",
)


def is_supported_url(url: str) -> bool:
    return any(h in url.lower() for h in SOCIAL_HOSTS)


async def download_media(url: str, audio_only: bool = False) -> Path:
    out_id = uuid.uuid4().hex
    out_template = str(DOWNLOAD_DIR / f"{out_id}.%(ext)s")
    if audio_only:
        ydl_opts = {
            "format": "bestaudio/best",
            "outtmpl": out_template,
            "postprocessors": [
                {
                    "key": "FFmpegExtractAudio",
                    "preferredcodec": "mp3",
                    "preferredquality": "192",
                }
            ],
            "quiet": True,
            "no_warnings": True,
            "noplaylist": True,
        }
    else:
        ydl_opts = {
            "format": "best[filesize<48M]/best[height<=720]/best",
            "outtmpl": out_template,
            "merge_output_format": "mp4",
            "quiet": True,
            "no_warnings": True,
            "noplaylist": True,
        }

    def _run():
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])

    await asyncio.get_event_loop().run_in_executor(None, _run)
    files = sorted(
        DOWNLOAD_DIR.glob(f"{out_id}.*"), key=lambda p: p.stat().st_mtime, reverse=True
    )
    if not files:
        raise RuntimeError("Yuklab olinmadi.")
    return files[0]


async def send_video_with_music_button(message: types.Message, path: Path, url: str):
    size_mb = path.stat().st_size / (1024 * 1024)
    if size_mb > 50:
        await message.answer(
            f"❌ Fayl juda katta ({size_mb:.1f} MB). Telegram cheklovi: 50 MB.",
            reply_markup=main_menu(message.from_user.id),
        )
        return
    async with aiofiles.open(path, "rb") as f:
        data = await f.read()
    bio = BytesIO(data)
    bio.name = path.name
    key = cache_url(url)
    await message.answer_video(
        bio,
        caption="🎬 Video tayyor",
        reply_markup=music_after_video_kb(key),
    )


async def send_audio_file(chat_id: int, path: Path):
    size_mb = path.stat().st_size / (1024 * 1024)
    if size_mb > 50:
        await bot.send_message(chat_id, f"❌ Audio juda katta ({size_mb:.1f} MB).")
        return
    async with aiofiles.open(path, "rb") as f:
        data = await f.read()
    bio = BytesIO(data)
    bio.name = path.name
    await bot.send_audio(chat_id, bio, caption="🎵 Musiqa tayyor")


# ---------- /start ----------
@dp.message_handler(commands=["start"], state="*")
async def cmd_start(message: types.Message, state: FSMContext):
    await state.finish()
    user_row = ensure_user(message.from_user)
    name = message.from_user.full_name
    admin_note = (
        "\n\n👑 <i>Siz admin sifatida belgilandingiz.</i>"
        if user_row["is_admin"]
        else ""
    )
    await message.answer(
        f"👋 Salom, <b>{name}</b>!\n\n"
        "Men ko'p funksiyali yordamchi botman:\n"
        "• 🎬 Video va 🎵 musiqa yuklash — <b>BEPUL</b> ♾ (YouTube, TikTok, Instagram, ...)\n"
        "• 🤖 AI yordamchi — <b>BEPUL</b> ♾\n"
        "• 📝 Referat / 📚 Kurs ishi / 📊 Slayd — har biri 1 marta bepul, keyingilari pullik\n\n"
        f"💡 Havola yuborsangiz, video chiqadi. Ostidagi tugma orqali musiqasini ham olishingiz mumkin."
        f"{admin_note}",
        reply_markup=main_menu(message.from_user.id),
    )


@dp.message_handler(commands=["help"], state="*")
async def cmd_help(message: types.Message):
    await message.answer(
        "<b>Buyruqlar:</b>\n"
        "/start — Asosiy menyu\n"
        "/cancel — Bekor qilish\n"
        "/myid — Telegram ID\n\n"
        "Tugmalardan tanlab foydalaning yoki to'g'ridan-to'g'ri havola yuboring.",
        reply_markup=main_menu(message.from_user.id),
    )


@dp.message_handler(commands=["myid"], state="*")
async def cmd_myid(message: types.Message):
    await message.answer(f"Sizning ID: <code>{message.from_user.id}</code>")


@dp.message_handler(commands=["cancel"], state="*")
@dp.message_handler(lambda m: m.text == "❌ Bekor qilish", state="*")
async def cmd_cancel(message: types.Message, state: FSMContext):
    await state.finish()
    await message.answer("Bekor qilindi.", reply_markup=main_menu(message.from_user.id))


@dp.message_handler(lambda m: m.text == "ℹ️ Yordam")
async def btn_help(message: types.Message):
    await cmd_help(message)


# ---------- Profile ----------
@dp.message_handler(lambda m: m.text == "👤 Profil")
async def btn_profile(message: types.Message):
    ensure_user(message.from_user)
    con = db()
    u = con.execute(
        "SELECT * FROM users WHERE user_id=?", (message.from_user.id,)
    ).fetchone()
    paid_count = con.execute(
        "SELECT COUNT(*) c FROM requests WHERE user_id=? AND status IN ('approved','done')",
        (message.from_user.id,),
    ).fetchone()["c"]
    con.close()
    role = "👑 Admin (cheksiz)" if u["is_admin"] else "👤 Foydalanuvchi"
    await message.answer(
        f"👤 <b>Profil</b>\n\n"
        f"Ism: {u['full_name']}\n"
        f"ID: <code>{u['user_id']}</code>\n"
        f"Holat: {role}\n\n"
        f"<b>Bepul foydalanish ({FREE_LIMIT} marta):</b>\n"
        f"📝 Referat: {u['essay_used']}\n"
        f"📚 Kurs ishi: {u['coursework_used']}\n"
        f"📊 Slayd: {u['slides_used']}\n"
        f"🎬 Video / 🎵 Musiqa: ♾ cheksiz bepul\n"
        f"🤖 AI yordamchi: ♾ cheksiz bepul\n\n"
        f"💳 Pullik buyurtmalar: {paid_count}",
        reply_markup=main_menu(message.from_user.id),
    )


# ---------- Student work entry points ----------
async def start_work(message: types.Message, state_obj: State, label: str):
    ensure_user(message.from_user)
    await state_obj.set()
    await message.answer(f"{label} mavzusini yozing:", reply_markup=cancel_kb())


@dp.message_handler(lambda m: m.text == "📝 Referat")
async def btn_essay(message: types.Message):
    await start_work(message, Form.waiting_essay_topic, "Referat")


@dp.message_handler(lambda m: m.text == "📚 Kurs ishi")
async def btn_coursework(message: types.Message):
    await start_work(message, Form.waiting_coursework_topic, "Kurs ishi")


@dp.message_handler(lambda m: m.text == "📊 Slayd (PPTX)")
async def btn_slides(message: types.Message):
    await start_work(message, Form.waiting_slides_topic, "Slayd")


async def after_topic(message: types.Message, state: FSMContext, rtype: str):
    topic = message.text.strip()
    if len(topic) < 3:
        await message.answer("Mavzu juda qisqa. Qayta yozing:")
        return
    await state.update_data(topic=topic)
    label, _ = REQUEST_LABELS[rtype]
    await message.answer(
        f"<b>{label}</b> uchun hajmni tanlang:\n\n"
        f"💡 Sizda <b>{FREE_LIMIT} marta bepul</b> imkoniyat bor — agar foydalanmagan bo'lsangiz, "
        f"istalgan hajmni tanlang, bepul tayyor bo'ladi.",
        reply_markup=pages_kb(rtype),
    )


@dp.message_handler(state=Form.waiting_essay_topic)
async def essay_topic(message: types.Message, state: FSMContext):
    await after_topic(message, state, "essay")


@dp.message_handler(state=Form.waiting_coursework_topic)
async def coursework_topic(message: types.Message, state: FSMContext):
    await after_topic(message, state, "coursework")


@dp.message_handler(state=Form.waiting_slides_topic)
async def slides_topic(message: types.Message, state: FSMContext):
    await after_topic(message, state, "slides")


# ---------- Size selection ----------
@dp.callback_query_handler(lambda c: c.data and c.data.startswith("sz:"), state="*")
async def cb_size(call: types.CallbackQuery, state: FSMContext):
    _, rtype, pages_s = call.data.split(":")
    pages = int(pages_s)
    pricing = dict(REQUEST_LABELS[rtype][1])
    price = pricing.get(pages)
    if price is None:
        await call.answer("Noto'g'ri tanlov", show_alert=True)
        return
    data = await state.get_data()
    topic = data.get("topic")
    if not topic:
        await call.answer("Mavzu topilmadi. Qaytadan boshlang.", show_alert=True)
        await state.finish()
        return
    await call.answer()
    user_id = call.from_user.id
    label = REQUEST_LABELS[rtype][0]

    # Free quota or admin → generate immediately
    if has_free_quota(user_id, rtype):
        await state.finish()
        await call.message.edit_reply_markup()
        await deliver_work(
            call.message.chat.id, user_id, rtype, topic, pages, free=True
        )
        if not is_admin(user_id):
            increment_usage(user_id, f"{rtype}_used")
        return

    # Otherwise — create paid request
    rid = create_request(user_id, rtype, topic, pages, price)
    await state.finish()
    await call.message.edit_reply_markup()
    await call.message.answer(
        f"💳 <b>To'lov</b>\n\n"
        f"📌 Buyurtma: <b>{label}</b> — <b>{topic}</b>\n"
        f"📄 Hajm: <b>{pages} bet</b>\n"
        f"💰 Narx: <b>{price:,} so'm</b>\n\n"
        f"To'lov uchun karta:\n"
        f"<code>{CARD_NUMBER}</code>\n"
        f"👤 {CARD_HOLDER}\n\n"
        f"To'lovni amalga oshirib, chek skrinshotini yuboring. "
        f"Admin tasdiqlagandan so'ng buyurtma sizga yuboriladi.\n\n"
        f"📋 Buyurtma ID: #{rid}",
        reply_markup=confirm_payment_kb(rid),
    )


@dp.callback_query_handler(lambda c: c.data and c.data.startswith("cancel:"), state="*")
async def cb_cancel_request(call: types.CallbackQuery):
    rid = int(call.data.split(":")[1])
    req = get_request(rid)
    if (
        req
        and req["user_id"] == call.from_user.id
        and req["status"] == "awaiting_payment"
    ):
        update_request(rid, status="cancelled")
    await call.answer("Bekor qilindi")
    await call.message.edit_reply_markup()
    await call.message.answer(
        "❌ Buyurtma bekor qilindi.", reply_markup=main_menu(call.from_user.id)
    )


@dp.callback_query_handler(lambda c: c.data and c.data.startswith("pay:"), state="*")
async def cb_pay(call: types.CallbackQuery, state: FSMContext):
    rid = int(call.data.split(":")[1])
    req = get_request(rid)
    if not req or req["user_id"] != call.from_user.id:
        await call.answer("Buyurtma topilmadi.", show_alert=True)
        return
    if req["status"] != "awaiting_payment":
        await call.answer(f"Holat: {req['status']}", show_alert=True)
        return
    await call.answer()
    await Form.waiting_receipt.set()
    await state.update_data(request_id=rid)
    await call.message.answer(
        f"📸 #{rid}-buyurtma uchun chek (skrinshot)ni rasm sifatida yuboring.",
        reply_markup=cancel_kb(),
    )


@dp.message_handler(state=Form.waiting_receipt, content_types=types.ContentType.PHOTO)
async def receive_receipt(message: types.Message, state: FSMContext):
    data = await state.get_data()
    rid = data.get("request_id")
    if not rid:
        await state.finish()
        await message.answer("Xato.", reply_markup=main_menu(message.from_user.id))
        return
    req = get_request(rid)
    if not req or req["user_id"] != message.from_user.id:
        await state.finish()
        await message.answer(
            "Buyurtma topilmadi.", reply_markup=main_menu(message.from_user.id)
        )
        return
    file_id = message.photo[-1].file_id
    update_request(rid, receipt_file_id=file_id, status="awaiting_review")
    await state.finish()

    admin_id = get_admin_id()
    if admin_id is None:
        await message.answer(
            "❌ Admin topilmadi.", reply_markup=main_menu(message.from_user.id)
        )
        return

    u = message.from_user
    label = REQUEST_LABELS[req["request_type"]][0]
    caption = (
        f"🧾 <b>Yangi to'lov cheki</b>\n\n"
        f"📋 #{rid} — {label}\n"
        f"📝 Mavzu: {req['topic']}\n"
        f"📄 Hajm: {req['pages']} bet\n"
        f"💰 Summa: {req['price']:,} so'm\n\n"
        f"👤 {u.full_name}\n"
        f"🆔 <code>{u.id}</code>\n"
        f"📛 @{u.username if u.username else '—'}\n"
        f"📅 {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )
    sent = await bot.send_photo(
        admin_id, file_id, caption=caption, reply_markup=admin_review_kb(rid)
    )
    update_request(rid, admin_msg_id=sent.message_id)

    await message.answer(
        "✅ Chek admin ko'rib chiqishga yuborildi.\n"
        "Tasdiqlangach, buyurtma sizga avtomatik yuboriladi.",
        reply_markup=main_menu(message.from_user.id),
    )


@dp.message_handler(state=Form.waiting_receipt, content_types=types.ContentType.ANY)
async def receipt_wrong(message: types.Message):
    await message.answer("Iltimos, faqat <b>rasm</b> yuboring (chek skrinshoti).")


# ---------- Admin review ----------
@dp.callback_query_handler(lambda c: c.data and c.data.startswith("adm:"))
async def cb_admin(call: types.CallbackQuery):
    if not is_admin(call.from_user.id):
        await call.answer("Faqat admin uchun.", show_alert=True)
        return
    _, action, rid_s = call.data.split(":")
    rid = int(rid_s)
    req = get_request(rid)
    if not req:
        await call.answer("Topilmadi.", show_alert=True)
        return
    if req["status"] != "awaiting_review":
        await call.answer(f"Holat: {req['status']}", show_alert=True)
        return

    if action == "approve":
        update_request(
            rid, status="approved", reviewed_at=datetime.utcnow().isoformat()
        )
        await call.answer("Tasdiqlandi. Buyurtma yuborilmoqda...")
        try:
            await call.message.edit_caption(
                (call.message.caption or "") + "\n\n✅ <b>QABUL QILINDI</b>",
                reply_markup=None,
            )
        except Exception:
            pass
        await bot.send_message(
            req["user_id"],
            f"✅ #{rid}-buyurtmangiz tasdiqlandi! Tayyorlanmoqda...",
        )
        try:
            await deliver_work(
                req["user_id"],
                req["user_id"],
                req["request_type"],
                req["topic"],
                req["pages"],
                free=False,
            )
            update_request(rid, status="done")
        except Exception as e:
            logger.exception("delivery failed")
            await bot.send_message(req["user_id"], f"❌ Yetkazib berishda xato: {e}")
    else:
        update_request(
            rid, status="rejected", reviewed_at=datetime.utcnow().isoformat()
        )
        try:
            await call.message.edit_caption(
                (call.message.caption or "") + "\n\n❌ <b>RAD ETILDI</b>",
                reply_markup=None,
            )
        except Exception:
            pass
        await call.answer("Rad etildi.")
        try:
            await bot.send_message(
                req["user_id"],
                f"❌ #{rid}-buyurtmangiz rad etildi.\n"
                "To'lov amalga oshganligini tekshiring va qayta urinib ko'ring.",
                reply_markup=main_menu(req["user_id"]),
            )
        except Exception:
            pass


# ---------- Generate & deliver ----------
async def deliver_work(
    chat_id: int, user_id: int, rtype: str, topic: str, pages: int, free: bool
):
    label = REQUEST_LABELS[rtype][0]
    text_label = "bepul" if free else "to'lovingiz uchun"
    if rtype == "slides":
        loading = await show_loading(chat_id, "Slayd tayyorlanmoqda")
        slides = await generate_slides_content(topic, n=pages)
        pptx = build_pptx(topic, slides)
        pptx.name = f"taqdimot_{topic[:30]}.pptx"
        await stop_loading(loading)
        await bot.send_document(
            chat_id,
            pptx,
            caption=f"📊 <b>{topic}</b> ({pages} slayd) — {text_label}",
            reply_markup=main_menu(user_id),
        )
        return

    loading = await show_loading(chat_id, f"{label} yozilmoqda")
    body = (
        await generate_essay(topic, pages)
        if rtype == "essay"
        else await generate_coursework(topic, pages)
    )
    docx = build_docx(f"{label}: {topic}", body)
    docx.name = f"{rtype}_{topic[:30]}.docx"
    await stop_loading(loading)
    await bot.send_document(
        chat_id,
        docx,
        caption=f"{label}: <b>{topic}</b> ({pages} bet) — {text_label}",
        reply_markup=main_menu(user_id),
    )


# ---------- Media (always free) ----------
@dp.message_handler(lambda m: m.text == "🎵 Musiqa/Video")
async def btn_media(message: types.Message):
    ensure_user(message.from_user)
    await message.answer(
        "Havolani yuboring (YouTube, TikTok, Instagram, Facebook, X/Twitter, SoundCloud va h.k.).\n"
        "Video chiqadi, ostida 🎵 <b>Musiqasini yuklab olish</b> tugmasi bo'ladi.",
        reply_markup=main_menu(message.from_user.id),
    )


async def handle_media_url(message: types.Message, url: str):
    if not is_supported_url(url):
        await message.answer(
            "Bu sayt qo'llab-quvvatlanmaydi.",
            reply_markup=main_menu(message.from_user.id),
        )
        return
    loading = await show_loading(message.chat.id, "Video yuklanmoqda")
    video_path: Path | None = None
    try:
        video_path = await download_media(url, audio_only=False)
        await stop_loading(loading)
        await send_video_with_music_button(message, video_path, url)
    except Exception as e:
        await stop_loading(loading)
        await message.answer(
            f"❌ Video yuklab bo'lmadi: {e}",
            reply_markup=main_menu(message.from_user.id),
        )
    finally:
        if video_path and video_path.exists():
            try:
                video_path.unlink()
            except Exception:
                pass


@dp.callback_query_handler(lambda c: c.data and c.data.startswith("audio:"))
async def cb_audio(call: types.CallbackQuery):
    key = call.data.split(":", 1)[1]
    url = URL_CACHE.get(key)
    if not url:
        await call.answer("Havola eskirdi. Qaytadan yuboring.", show_alert=True)
        return
    await call.answer("Musiqa yuklanmoqda...")
    loading = await show_loading(call.message.chat.id, "Musiqa ajratilmoqda")
    audio_path: Path | None = None
    try:
        audio_path = await download_media(url, audio_only=True)
        await stop_loading(loading)
        await send_audio_file(call.message.chat.id, audio_path)
    except Exception:
        await stop_loading(loading)
        await bot.send_message(call.message.chat.id, "🎵 Musiqa topilmadi.")
    finally:
        if audio_path and audio_path.exists():
            try:
                audio_path.unlink()
            except Exception:
                pass


# ---------- AI chat (always free) ----------
@dp.message_handler(lambda m: m.text == "🤖 AI yordamchi")
async def btn_ai(message: types.Message, state: FSMContext):
    ensure_user(message.from_user)
    await Form.chatting_ai.set()
    await state.update_data(history=[])
    await message.answer(
        "🤖 AI yordamchi rejimi yoqildi (cheksiz bepul ♾).\n"
        "Savollaringizni yuboring. Chiqish — ❌ Bekor qilish.",
        reply_markup=cancel_kb(),
    )


@dp.message_handler(state=Form.chatting_ai, content_types=types.ContentType.TEXT)
async def ai_chat(message: types.Message, state: FSMContext):
    if message.text == "❌ Bekor qilish":
        await state.finish()
        await message.answer(
            "Bekor qilindi.", reply_markup=main_menu(message.from_user.id)
        )
        return
    data = await state.get_data()
    history = data.get("history", [])
    history.append({"role": "user", "content": message.text})
    loading = await show_loading(message.chat.id, "AI o'ylayapti")
    if not ai_client:
        await stop_loading(loading)
        await message.answer("AI sozlanmagan.")
        return
    try:
        resp = await ai_client.chat.completions.create(
            model="gpt-5.4",
            max_completion_tokens=8192,
            messages=[
                {
                    "role": "system",
                    "content": "Siz foydali AI yordamchi. O'zbek tilida qisqa va aniq javob bering.",
                },
                *history[-10:],
            ],
        )
        answer = resp.choices[0].message.content or "..."
    except Exception as e:
        answer = f"Xato: {e}"
    history.append({"role": "assistant", "content": answer})
    await state.update_data(history=history)
    await stop_loading(loading)
    await message.answer(answer)


# ---------- Admin panel ----------
@dp.message_handler(lambda m: m.text == "🛠 Admin panel")
async def btn_admin(message: types.Message):
    if not is_admin(message.from_user.id):
        return
    con = db()
    users_count = con.execute("SELECT COUNT(*) c FROM users").fetchone()["c"]
    pending = con.execute(
        "SELECT COUNT(*) c FROM requests WHERE status='awaiting_review'"
    ).fetchone()["c"]
    approved = con.execute(
        "SELECT COUNT(*) c FROM requests WHERE status IN ('approved','done')"
    ).fetchone()["c"]
    revenue = con.execute(
        "SELECT COALESCE(SUM(price),0) s FROM requests WHERE status IN ('approved','done')"
    ).fetchone()["s"]
    con.close()
    await message.answer(
        f"🛠 <b>Admin panel</b>\n\n"
        f"👥 Foydalanuvchilar: {users_count}\n"
        f"⏳ Tekshirilmagan to'lovlar: {pending}\n"
        f"✅ Bajarilgan buyurtmalar: {approved}\n"
        f"💰 Umumiy daromad: {revenue:,} so'm",
        reply_markup=main_menu(message.from_user.id),
    )


# ---------- Fallback ----------
@dp.message_handler(content_types=types.ContentType.TEXT, state=None)
async def fallback(message: types.Message):
    ensure_user(message.from_user)
    text = message.text.strip()
    url_match = re.search(r"https?://\S+", text)
    if url_match and is_supported_url(url_match.group(0)):
        await handle_media_url(message, url_match.group(0))
        return
    # Otherwise — free AI answer
    loading = await show_loading(message.chat.id, "Javob tayyorlanmoqda")
    answer = await ai_text(text)
    await stop_loading(loading)
    await message.answer(answer, reply_markup=main_menu(message.from_user.id))


if __name__ == "__main__":
    db_init()
    keep_alive()
    logger.info("Bot ishga tushdi.")
    executor.start_polling(dp, skip_updates=True)
